# -*- coding: utf-8 -*-
"""
增强版文件加载模块
支持多种文件格式：
- PDF（纯文本和图片PDF，带OCR支持）
- DOC/DOCX
- PPT/PPTX
- 图片格式（PNG, JPG, JPEG, BMP, TIFF等）
- TXT文本文件
- MD markdown文件
"""

import os
import platform
import pytesseract
from PIL import Image
from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from docx import Document
from pptx import Presentation

_tesseract_path = os.getenv("TESSERACT_PATH", "")
if _tesseract_path:
    pytesseract.pytesseract.tesseract_cmd = _tesseract_path
elif platform.system() == "Windows":
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text_from_image(image_path):
    """
    使用OCR从图片中提取文本
    :param image_path: 图片文件路径
    :return: 提取的文本字符串
    """
    try:
        img = Image.open(image_path)
        # 尝试转换为灰度图像以提高OCR准确性
        img = img.convert('L')
        text = pytesseract.image_to_string(img, lang='chi_sim')
        return text
    except Exception as e:
        raise RuntimeError(f"OCR识别失败：{image_path}，错误原因：{str(e)}")

def extract_text_from_pdf_with_ocr(pdf_path):
    """
    使用OCR从图片PDF中提取文本
    :param pdf_path: PDF文件路径
    :return: 提取的文本字符串
    """
    try:
        from pdf2image import convert_from_path
        
        # 将PDF转换为图片
        pages = convert_from_path(pdf_path)
        
        # 对每一页进行OCR
        text = ""
        for page in pages:
            page_text = pytesseract.image_to_string(page, lang='chi_sim')
            if page_text:
                text += page_text + "\n"
        
        return text
    except ImportError:
        raise ImportError("需要安装pdf2image库：pip install pdf2image")
    except Exception as e:
        raise RuntimeError(f"PDF OCR识别失败：{pdf_path}，错误原因：{str(e)}")

def load_pdf(pdf_path: str, use_ocr: bool = False) -> str:
    """
    读取PDF文件，支持纯文本和图片PDF（通过OCR）
    :param pdf_path: PDF文件路径
    :param use_ocr: 是否使用OCR（当纯文本提取失败时自动尝试）
    :return: 文本字符串
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF文件不存在：{pdf_path}")

    if not pdf_path.lower().endswith('.pdf'):
        raise ValueError(f"文件不是PDF格式：{pdf_path}")

    try:
        reader = PdfReader(pdf_path)

        if reader.is_encrypted:
            raise PdfReadError(f"PDF文件已加密，无法读取：{pdf_path}")

        # 先尝试纯文本提取
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        # 如果启用了OCR且纯文本提取结果很少或为空，则尝试OCR
        if use_ocr and len(text.strip()) < 10:
                # 可能是图片PDF，尝试OCR
                ocr_text = extract_text_from_pdf_with_ocr(pdf_path)
                if len(ocr_text.strip()) > len(text.strip()):
                    text = ocr_text

        return text

    except PdfReadError as e:
        raise PdfReadError(f"解析PDF失败：{pdf_path}，错误原因：{str(e)}")

def load_docx(docx_path: str) -> str:
    """
    读取DOCX文件
    :param docx_path: DOCX文件路径
    :return: 文本字符串
    """
    if not os.path.exists(docx_path):
        raise FileNotFoundError(f"DOCX文件不存在：{docx_path}")

    if not docx_path.lower().endswith('.docx'):
        raise ValueError(f"文件不是DOCX格式：{docx_path}")

    try:
        doc = Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        raise RuntimeError(f"解析DOCX失败：{docx_path}，错误原因：{str(e)}")

def load_ppt(ppt_path: str) -> str:
    """
    读取PPT/PPTX文件
    :param ppt_path: PPT/PPTX文件路径
    :return: 文本字符串
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PPT文件不存在：{ppt_path}")

    if not (ppt_path.lower().endswith('.ppt') or ppt_path.lower().endswith('.pptx')):
        raise ValueError(f"文件不是PPT格式：{ppt_path}")

    try:
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise RuntimeError(f"解析PPT失败：{ppt_path}，错误原因：{str(e)}")

def load_image(image_path: str) -> str:
    """
    使用OCR从图片文件中提取文本
    :param image_path: 图片文件路径
    :return: 文本字符串
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"图片文件不存在：{image_path}")

    supported_formats = ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif')
    if not any(image_path.lower().endswith(fmt) for fmt in supported_formats):
        raise ValueError(f"不支持的图片格式：{image_path}")

    return extract_text_from_image(image_path)

def load_txt(txt_path: str) -> str:
    """
    读取TXT文本文件
    :param txt_path: TXT文件路径
    :return: 文本字符串
    """
    if not os.path.exists(txt_path):
        raise FileNotFoundError(f"TXT文件不存在：{txt_path}")

    if not txt_path.lower().endswith('.txt'):
        raise ValueError(f"文件不是TXT格式：{txt_path}")

    try:
        # 尝试多种编码读取
        encodings = ['utf-8', 'gbk', 'gb2312', 'big5']
        for encoding in encodings:
            try:
                with open(txt_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise RuntimeError(f"无法解码文件：{txt_path}")
    except Exception as e:
        raise RuntimeError(f"读取TXT失败：{txt_path}，错误原因：{str(e)}")

def load_md(md_path: str) -> str:
    """
    读取MD Markdown文件
    :param md_path: MD文件路径
    :return: 文本字符串（保留原始内容）
    """
    if not os.path.exists(md_path):
        raise FileNotFoundError(f"MD文件不存在：{md_path}")

    if not md_path.lower().endswith('.md'):
        raise ValueError(f"文件不是MD格式：{md_path}")

    try:
        with open(md_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        raise RuntimeError(f"读取MD失败：{md_path}，错误原因：{str(e)}")

def load_file(file_path: str, use_ocr: bool = True) -> str:
    """
    通用文件加载函数，根据文件扩展名自动选择合适的加载方式
    :param file_path: 文件路径
    :param use_ocr: 是否启用OCR（适用于图片PDF和图片文件）
    :return: 提取的文本字符串
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在：{file_path}")

    # 获取文件扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    # 根据扩展名选择加载方式
    if ext == '.pdf':
        return load_pdf(file_path, use_ocr)
    elif ext == '.docx':
        return load_docx(file_path)
    elif ext in ('.ppt', '.pptx'):
        return load_ppt(file_path)
    elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'):
        return load_image(file_path)
    elif ext == '.txt':
        return load_txt(file_path)
    elif ext == '.md':
        return load_md(file_path)
    else:
        raise ValueError(f"不支持的文件格式：{ext}")

def get_supported_formats() -> list:
    """
    获取支持的文件格式列表
    :return: 支持的文件格式列表
    """
    return [
        '.pdf',      # PDF文件（支持纯文本和图片PDF）
        '.docx',     # Word文档
        '.ppt',      # PowerPoint演示文稿
        '.pptx',     # PowerPoint演示文稿
        '.png',      # PNG图片
        '.jpg',      # JPG图片
        '.jpeg',     # JPEG图片
        '.bmp',      # BMP图片
        '.tiff',     # TIFF图片
        '.gif',      # GIF图片
        '.txt',      # 文本文件
        '.md'        # Markdown文件
    ]